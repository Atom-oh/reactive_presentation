#!/usr/bin/env python3
"""
Extract specific slides from a PPTX file by copying the entire file
and then deleting unwanted slides. This preserves layouts, masters, and themes.

15-minute Checkpointless Training session (slides 23-36 + 50-52 from original)
"""

import shutil
from pptx import Presentation

SRC = '/home/ec2-user/reactive_presentation/AIM3338_[NEW-LAUNCH]-SageMaker-HyperPod-Checkpointless-and-elastic-training-for-AI-models.pptx'
DST = '/home/ec2-user/reactive_presentation/checkpointless-training/AIM3338_Checkpointless_15min.pptx'

# Slides to keep (0-indexed):
# - Slides 22-35: Checkpointless training section (original slides 23-36)
# - Slides 49-51: Takeaways and Resources (original slides 50-52)
KEEP_INDICES = set(list(range(22, 36)) + list(range(49, 52)))

def get_slide_title(slide):
    """Extract the first text from a slide for identification."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text[:80] + ('...' if len(text) > 80 else '')
    return "(no text)"

def delete_slide(prs, index):
    """Delete a slide at the given index using XML manipulation."""
    slide = prs.slides[index]
    slide_id = slide.slide_id

    prs_part = prs.part
    prs_xml = prs_part._element

    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    }

    sldIdLst = prs_xml.find('.//p:sldIdLst', nsmap)

    for sldId in sldIdLst.findall('p:sldId', nsmap):
        if int(sldId.get('id')) == slide_id:
            rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            sldIdLst.remove(sldId)
            if rId:
                prs_part.drop_rel(rId)
            break

def main():
    print(f"Source: {SRC}")
    print(f"Destination: {DST}")
    print(f"Keeping slide indices: {sorted(KEEP_INDICES)}")
    print(f"Total slides to keep: {len(KEEP_INDICES)}")
    print()

    # Step 1: Copy the entire source file
    print("Copying source file...")
    shutil.copy2(SRC, DST)

    # Step 2: Open the copy and identify slides to remove
    prs = Presentation(DST)
    total_slides = len(prs.slides)
    print(f"Original presentation has {total_slides} slides")

    # Validate keep indices
    invalid = [i for i in KEEP_INDICES if i >= total_slides]
    if invalid:
        print(f"WARNING: Invalid slide indices (out of range): {invalid}")

    # Print what we're keeping
    print("\nSlides to KEEP:")
    for i in sorted(KEEP_INDICES):
        if i < total_slides:
            slide = prs.slides[i]
            title = get_slide_title(slide)
            print(f"  [{i:2d}] {title}")

    # Calculate indices to remove (in reverse order to avoid index shifting)
    remove_indices = sorted([i for i in range(total_slides) if i not in KEEP_INDICES], reverse=True)
    print(f"\nRemoving {len(remove_indices)} slides...")

    # Step 3: Remove slides in reverse order (highest index first)
    for idx in remove_indices:
        delete_slide(prs, idx)

    # Save the modified presentation
    print("Saving modified presentation...")
    prs.save(DST)

    # Verify the result
    print("\n" + "="*60)
    print("VERIFICATION")
    print("="*60)

    prs_verify = Presentation(DST)
    final_count = len(prs_verify.slides)
    print(f"Final slide count: {final_count}")
    print(f"Expected: {len(KEEP_INDICES)}")

    if final_count == len(KEEP_INDICES):
        print("SUCCESS: Slide count matches!")
    else:
        print(f"WARNING: Expected {len(KEEP_INDICES)}, got {final_count}")

    print("\nExtracted slides:")
    for i, slide in enumerate(prs_verify.slides):
        title = get_slide_title(slide)
        print(f"  [{i:2d}] {title}")

    print(f"\nOutput saved to: {DST}")

if __name__ == '__main__':
    main()
